# -*- coding: utf-8 -*-
"""Populate the ISRO BAH 2026 idea-submission template with real project content."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SRC = "[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx"
DST = "[Pub] ISRO BAH 2026 _ Idea Submission Template.pptx"

NAVY = RGBColor(0x0D, 0x2A, 0x4A)
DARK = RGBColor(0x22, 0x22, 0x22)
ORANGE = RGBColor(0xE6, 0x51, 0x00)
GREY = RGBColor(0x55, 0x55, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

REPO = "https://github.com/Navn2025/DPSR"

prs = Presentation(SRC)
slides = prs.slides


def clear_shape_text(shape):
    tf = shape.text_frame
    for p in list(tf.paragraphs):
        for r in list(p.runs):
            r.text = ""
    # remove all but first paragraph
    txBody = tf._txBody
    ps = txBody.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p')
    for p in ps[1:]:
        txBody.remove(p)
    tf.paragraphs[0].text = ""


def set_box(shape, left=None, top=None, width=None, height=None):
    if left is not None:
        shape.left = Inches(left)
    if top is not None:
        shape.top = Inches(top)
    if width is not None:
        shape.width = Inches(width)
    if height is not None:
        shape.height = Inches(height)


def add_title(slide, text, top=0.72, left=0.3, width=9.4, size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.42))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = "Calibri"
    return box


def add_bullets(slide, items, left, top, width, height, size=12.5, header_size=13.5):
    """items: list of (kind, text) where kind in {'h','b','n'} header/bullet/note"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for kind, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run()
        if kind == "h":
            r.text = text
            r.font.bold = True
            r.font.size = Pt(header_size)
            r.font.color.rgb = NAVY
            p.space_before = Pt(8)
            p.space_after = Pt(2)
        elif kind == "b":
            r.text = "•  " + text
            r.font.size = Pt(size)
            r.font.color.rgb = DARK
            p.space_after = Pt(3)
        else:  # note / plain
            r.text = text
            r.font.size = Pt(size - 0.5)
            r.font.italic = True
            r.font.color.rgb = GREY
            p.space_after = Pt(3)
    return box


def add_footer_link(slide, text=f"Repo: {REPO}", top=5.05, color=GREY, left=0.3, width=9.4, size=10, align=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.italic = True
    return box


def add_pic_fit(slide, path, left, top, max_w, max_h):
    from PIL import Image
    im = Image.open(path)
    ratio = im.size[0] / im.size[1]
    w = max_w
    h = w / ratio
    if h > max_h:
        h = max_h
        w = h * ratio
    pic = slide.shapes.add_picture(path, Inches(left + (max_w - w) / 2), Inches(top), width=Inches(w), height=Inches(h))
    return pic, w, h


# =============================================================================
# SLIDE 1 — Title
# =============================================================================
s1 = slides[0]
texts = {shp.shape_id: shp for shp in s1.shapes if shp.has_text_frame}
for shp in s1.shapes:
    if not shp.has_text_frame:
        continue
    t = shp.text_frame.text.strip()
    if t.startswith("Team Name"):
        tf = shp.text_frame
        tf.paragraphs[0].runs[0].text = "Team Name :  "
        r2 = tf.paragraphs[0].add_run()
        r2.text = "[ FILL IN ]"
        r2.font.color.rgb = ORANGE
        r2.font.bold = True
    elif t.startswith("Problem Statement"):
        tf = shp.text_frame
        tf.paragraphs[0].runs[0].text = "Problem Statement :  "
        r2 = tf.paragraphs[0].add_run()
        r2.text = ("PS-8: Detection and Characterization of Subsurface Ice in Lunar South Polar "
                   "Regions Using Chandrayaan-2 Radar and Imagery Data for Landing Site and "
                   "Rover Traverse Planning")
        r2.font.size = Pt(11)
    elif t.startswith("Team Leader"):
        tf = shp.text_frame
        tf.paragraphs[0].runs[0].text = "Team Leader Name :  "
        r2 = tf.paragraphs[0].add_run()
        r2.text = "[ FILL IN ]"
        r2.font.color.rgb = ORANGE
        r2.font.bold = True

# =============================================================================
# SLIDE 3 — Opportunity
# =============================================================================
s3 = slides[2]
for shp in list(s3.shapes):
    if shp.has_text_frame and "Opportunity should be" in shp.text_frame.text:
        shp._element.getparent().remove(shp._element)

add_title(s3, "Opportunity — Scientific Differentiation & USP")
add_bullets(s3, [
    ("h", "How different from existing ideas"),
    ("b", "Most PSR/ice-mapping work publishes ONE evidence layer (shadow OR CPR OR temperature)."),
    ("b", "We fuse THREE independent signatures into one explainable, per-pixel Ice Confidence Score — no black-box ML."),
    ("b", "We implement PS-8's exact criterion, CPR>1 AND DOP<0.13, from the full Stokes derivation."),
    ("h", "How it solves the problem statement"),
    ("b", "PSR -> DPSR via curvature-corrected ray-casting on the LOLA DEM (dpsr/)"),
    ("b", "CPR + DOP computed directly from CH2 DFSAR SLI/GRI (cpr/, dop/)"),
    ("b", "CPR cross-validated pixel-by-pixel vs. ISRO's own DFSAR mosaic (validation/)"),
    ("b", "Fused with LRO Diviner thermal data into one Ice Confidence Map (diviner/)"),
    ("h", "USP"),
    ("b", "Validates against ISRO's own official DFSAR product, not just NASA refs"),
    ("b", "GPU-accelerated (CUDA/Numba) — full south-pole DPSR ray-tracing in minutes"),
], left=0.3, top=1.15, width=5.7, height=3.75, size=10.5)

add_pic_fit(s3, "temp/diagrams/key_equations.png", left=6.15, top=1.15, max_w=3.35, max_h=3.75)
add_footer_link(s3, top=5.10)

# =============================================================================
# SLIDE 4 — Features
# =============================================================================
s4 = slides[3]
for shp in list(s4.shapes):
    if shp.has_text_frame and "List of features" in shp.text_frame.text:
        shp._element.getparent().remove(shp._element)

add_title(s4, "Key Features & Measured Results")
add_bullets(s4, [
    ("h", "Implemented, end-to-end (this phase)"),
    ("b", "DPSR extraction — curvature-corrected ray-casting, Eq. A4, CPU/GPU (dpsr/)"),
    ("b", "CPR from CH2 DFSAR SLI + GRI — 2 formulas benchmarked (cpr/, cpr_gri/)"),
    ("b", "DOP from full Stokes derivation, no simplifying assumptions (dop/)"),
    ("b", "Automated validation vs. official ISRO CPR mosaic: r, SSIM, RMSE (validation/)"),
    ("b", "LRO Diviner fusion — Tmean, ZIT, Pump cold-trap efficiency (diviner/)"),
    ("b", "8-band physics-weighted Ice Confidence Map, full south pole, 230M px"),
    ("h", "Measured numbers"),
    ("b", "PSR = 10.79% of DEM area  |  DPSR = 0.0079% of DEM area"),
    ("b", "Faustini ice-candidates (CPR>1 & DOP<0.13): 15,153 px, consistent 8-12.6% across 9 independent azimuth strips"),
], left=0.3, top=1.18, width=5.5, height=3.85, size=11)

_, w, h = add_pic_fit(s4, "images-final/diviner/Ice_Confidence_Map.png", left=5.85, top=1.18, max_w=3.6, max_h=3.35)
cap = s4.shapes.add_textbox(Inches(5.85), Inches(1.18 + h + 0.05), Inches(3.6), Inches(0.5))
p = cap.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Ice Confidence Map — full LOLA south-pole grid"
r.font.size = Pt(9.5)
r.font.italic = True
r.font.color.rgb = GREY
add_footer_link(s4, top=5.08)

# =============================================================================
# SLIDE 5 — Process Flow
# =============================================================================
s5 = slides[4]
for shp in list(s5.shapes):
    if shp.has_text_frame and "Process flow diagram" in shp.text_frame.text:
        shp._element.getparent().remove(shp._element)

add_title(s5, "End-to-End Process Flow")
sub = s5.shapes.add_textbox(Inches(0.3), Inches(1.15), Inches(9.4), Inches(0.35))
p = sub.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Data -> per-sensor processing -> validation -> physics fusion -> ice-candidate output"
r.font.size = Pt(11)
r.font.italic = True
r.font.color.rgb = GREY

_, w, h = add_pic_fit(s5, "temp/diagrams/pipeline_flow.png", left=0.2, top=1.55, max_w=9.6, max_h=3.15)
legend = s5.shapes.add_textbox(Inches(0.3), Inches(1.55 + h + 0.08), Inches(9.4), Inches(0.4))
p = legend.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Grey = input data   Green = shadow modelling   Blue = radar polarimetry   Purple = fusion   Orange dashed = planned (Sprint 2-5, not yet executed)"
r.font.size = Pt(9.5)
r.font.color.rgb = GREY
add_footer_link(s5, top=5.08)

# =============================================================================
# SLIDE 6 — Wireframes -> Current Output Surface (honest, practical)
# =============================================================================
s6 = slides[5]
for shp in list(s6.shapes):
    if shp.has_text_frame and "Wireframes" in shp.text_frame.text:
        shp._element.getparent().remove(shp._element)

add_title(s6, "Current Output Surface  (Wireframes: not applicable this phase)")
sub = s6.shapes.add_textbox(Inches(0.3), Inches(1.15), Inches(9.4), Inches(0.75))
tf = sub.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
r = p.add_run()
r.text = ("No interactive dashboard has been built yet. The real, current output is a set of "
          "georeferenced GeoTIFF map layers, validated pixel-by-pixel against ISRO's own "
          "Chandrayaan-2 DFSAR CPR mosaic (Putrevu et al. 2023):")
r.font.size = Pt(11.5)
r.font.color.rgb = DARK

_, w, h = add_pic_fit(s6, "images-final/validation/validation_maps.png", left=0.9, top=1.95, max_w=8.2, max_h=2.85)
cap = s6.shapes.add_textbox(Inches(0.9), Inches(1.95 + h + 0.05), Inches(8.2), Inches(0.35))
p = cap.text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Our computed CPR (left) vs. official DFSAR mosaic (right) — SSIM = 0.944, same grid & colour scale"
r.font.size = Pt(9.5)
r.font.italic = True
r.font.color.rgb = GREY
add_footer_link(s6, top=5.08)

# =============================================================================
# SLIDE 7 — Architecture
# =============================================================================
s7 = slides[6]
for shp in list(s7.shapes):
    if shp.has_text_frame and "Architecture diagram" in shp.text_frame.text:
        shp._element.getparent().remove(shp._element)

add_title(s7, "System Architecture")
_, w, h = add_pic_fit(s7, "temp/diagrams/architecture_2row.png", left=0.2, top=1.15, max_w=9.6, max_h=3.55)
legend = s7.shapes.add_textbox(Inches(0.3), Inches(1.15 + h + 0.08), Inches(9.4), Inches(0.5))
p = legend.text_frame.paragraphs[0]
r = p.add_run()
r.text = "Solid boxes = implemented and executed. Dashed orange = designed in code (function signatures + formulas scoped) but not yet executed — Sprint 2-5 roadmap."
r.font.size = Pt(9.5)
r.font.color.rgb = GREY
add_footer_link(s7, top=5.08)

# =============================================================================
# SLIDE 8 — Technologies
# =============================================================================
s8 = slides[7]
for shp in list(s8.shapes):
    if shp.has_text_frame and "Technologies to be used" in shp.text_frame.text:
        shp._element.getparent().remove(shp._element)

add_title(s8, "Technologies, Data & Scientific References")
add_bullets(s8, [
    ("h", "Software & Libraries"),
    ("b", "Python 3.11"),
    ("b", "Rasterio / GDAL, GeoPandas — geospatial I/O & reprojection"),
    ("b", "Numba — JIT CPU parallel + CUDA GPU kernels"),
    ("b", "NumPy, SciPy — array math, Pearson/Spearman stats"),
    ("b", "scikit-image — SSIM"),
    ("b", "Matplotlib — all visualisation"),
    ("b", "QGIS — manual spot-check validation"),
    ("b", "Git — version control (repo below)"),
], left=0.3, top=1.18, width=4.5, height=3.85, size=11)

add_bullets(s8, [
    ("h", "Data Sources"),
    ("b", "LOLA DEM + PSR shapefile — PDS Geosciences Node"),
    ("b", "CH2 DFSAR SLI/GRI/SRI — ISRO ISDA, pradan.issdc.gov.in"),
    ("b", "LRO Diviner Tmean/ZIT/Pump — NASA PDS"),
    ("h", "Key References"),
    ("b", "Putrevu et al. (2023) JGR Planets — official CPR mosaic"),
    ("b", "O'Brien & Byrne (2022) PSJ 3:258 — DPSR ray-tracing"),
    ("b", "Nozette (1996) Science; Campbell (2006) Nature — CPR/ice"),
    ("b", "van Zyl & Kim (2011) — SAR Polarimetry / DOP"),
    ("b", "Paige (2010) Science; Hayne (2015) Icarus — thermal stability"),
    ("b", "Schorghofer (2014) ApJ — cold-trap efficiency"),
], left=5.0, top=1.18, width=4.7, height=3.85, size=11)
add_footer_link(s8, top=5.08)

# =============================================================================
# SLIDE 9 — Cost / Results snapshot
# =============================================================================
s9 = slides[8]
for shp in list(s9.shapes):
    if shp.has_text_frame and "Estimated implementation cost" in shp.text_frame.text:
        shp._element.getparent().remove(shp._element)

add_title(s9, "Results Snapshot & Estimated Implementation Cost")
add_bullets(s9, [
    ("h", "Compute footprint (measured)"),
    ("b", "Full annual-illumination DPSR ray-tracing: ~30-90 min GPU vs 2-4 hr CPU-only"),
    ("b", "Feature_Stack.tif: 5.2 GB (9 bands, 230M px)  |  Ice_Confidence_Map.tif: 992 MB"),
    ("b", "Single workstation + optional NVIDIA GPU — no specialised/mission hardware"),
    ("h", "Path to production (~3-4 person-months)"),
    ("b", "Fix 2 known integration bugs (DOP-to-fusion path, correlation matrix)"),
    ("b", "Execute Sprint 2-5 roadmap already scoped in code: ice-probability map -> landing-site selection -> A* rover path planning -> ice-volume estimate (top 5 m)"),
], left=0.3, top=1.18, width=4.9, height=3.85, size=11)

_, w, h = add_pic_fit(s9, "temp/diagrams/stats_cpr_validation.png", left=5.35, top=1.3, max_w=4.15, max_h=3.0)
cap = s9.shapes.add_textbox(Inches(5.35), Inches(1.3 + h + 0.05), Inches(4.15), Inches(0.55))
tf = cap.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "3 independently computed CPR products, validated vs. ISRO's official DFSAR mosaic"
r.font.size = Pt(9.5)
r.font.italic = True
r.font.color.rgb = GREY
add_footer_link(s9, top=5.08)

# =============================================================================
# SLIDE 10 — closing (design already has "THANK YOU"; just add repo footer)
# =============================================================================
s10 = slides[9]
box = s10.shapes.add_textbox(Inches(0.4), Inches(4.72), Inches(9.0), Inches(0.35))
p = box.text_frame.paragraphs[0]
r = p.add_run()
r.text = f"{REPO}   |   Team: [ FILL IN ]"
r.font.size = Pt(12)
r.font.bold = True
r.font.color.rgb = RGBColor(0xFF, 0xA5, 0x4D)

prs.save(DST)
print("Saved:", DST)
